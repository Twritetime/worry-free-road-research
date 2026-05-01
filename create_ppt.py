from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE_DARK = RGBColor(0, 82, 147)
BLUE_LIGHT = RGBColor(0, 122, 194)
BLUE_BG = RGBColor(230, 240, 250)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(100, 100, 100)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_DARK
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(54)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        tf.paragraphs[0].text = subtitle
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, two_column=False):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE_DARK
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    if two_column and len(content_items) >= 2:
        left_items = content_items[0]
        right_items = content_items[1]

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            if i == 0:
                tf.paragraphs[0].text = "  " + item
                tf.paragraphs[0].font.size = Pt(22)
                tf.paragraphs[0].font.color.rgb = BLACK
                tf.paragraphs[0].font.bold = (i == 0)
            else:
                p = tf.add_paragraph()
                p.text = "  " + item
                p.font.size = Pt(20)
                p.font.color.rgb = GRAY

        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            if i == 0:
                tf.paragraphs[0].text = "  " + item
                tf.paragraphs[0].font.size = Pt(22)
                tf.paragraphs[0].font.color.rgb = BLACK
            else:
                p = tf.add_paragraph()
                p.text = "  " + item
                p.font.size = Pt(20)
                p.font.color.rgb = GRAY
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(content_items):
            if i == 0:
                tf.paragraphs[0].text = "  " + item
                tf.paragraphs[0].font.size = Pt(24)
                tf.paragraphs[0].font.color.rgb = BLACK
                tf.paragraphs[0].font.bold = True
            else:
                p = tf.add_paragraph()
                p.text = "  " + item
                p.font.size = Pt(20)
                p.font.color.rgb = GRAY

    return slide

def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_LIGHT
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(48)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

slides_data = []

slides_data.append({"type": "title", "title": "研路无忧", "subtitle": "考研服务平台\n\n答辩人：吴涛\n软工223\n指导老师：王华"})

slides_data.append({"type": "content", "title": "目录", "content": [
    "一、项目背景与简介",
    "二、需求分析",
    "三、技术架构",
    "四、数据库设计",
    "五、功能模块展示",
    "六、系统测试与总结"
]})

slides_data.append({"type": "section", "title": "一、项目背景与简介"})

slides_data.append({"type": "content", "title": "项目背景", "content": [
    "  随着考研人数逐年增长，考研学生面临诸多挑战：",
    "",
    "  • 信息分散：考研资讯、备考资料分布在不同平台，查找困难",
    "  • 资料获取难：优质学习资料难以获取，价格不透明",
    "  • 交流匮乏：缺乏专业的考研交流社区",
    "  • 服务单一：现有平台功能单一，无法满足一站式需求",
    "",
    "  因此，我们开发了'研路无忧'考研服务平台，旨在为考研学生提供",
    "  便捷、高效的一站式服务。"
]})

slides_data.append({"type": "content", "title": "项目简介", "content": [
    "  项目名称：研路无忧（YanLuWuYou）",
    "",
    "  核心功能：",
    "  • 资讯中心：实时更新的考研动态与新闻（爬虫自动获取）",
    "  • 资料商城：学习资料浏览、购买、支付（支付宝集成）",
    "  • 交流论坛：发帖、评论、点赞收藏",
    "  • 备考指南：各科目备考建议与经验分享",
    "  • 管理后台：用户、帖子、资料、订单等统一管理",
    "",
    "  技术栈：Spring Boot + Vue 3 + MyBatis Plus + MySQL"
]})

slides_data.append({"type": "section", "title": "二、需求分析"})

slides_data.append({"type": "content", "title": "功能需求分析", "content": [
    "  1. 用户系统",
    "     用户注册、登录、个人信息维护、收货地址管理",
    "",
    "  2. 资讯中心",
    "     考研新闻自动爬取、分类展示、详情查看",
    "",
    "  3. 资料商城",
    "     资料浏览、搜索、加入购物车、下单购买、支付宝支付",
    "",
    "  4. 交流论坛",
    "     发布帖子、发表评论、点赞、收藏",
    "",
    "  5. 备考指南",
    "     备考经验文章分类展示、详细内容浏览",
    "",
    "  6. 反馈建议",
    "     用户提交反馈、管理员处理反馈"
]})

slides_data.append({"type": "content", "title": "非功能需求", "content": [
    "  性能需求",
    "  • 系统响应时间 < 2秒",
    "  • 支持并发用户数 > 100",
    "",
    "  安全需求",
    "  • 用户密码加密存储",
    "  • 接口身份认证与授权",
    "  • 防止SQL注入、XSS攻击",
    "",
    "  可用性需求",
    "  • 系统可用性 > 99%",
    "  • 界面简洁美观，操作便捷",
    "",
    "  可扩展性需求",
    "  • 模块化设计，便于功能扩展",
    "  • 数据库结构优化，便于数据扩展"
]})

slides_data.append({"type": "section", "title": "三、技术架构"})

slides_data.append({"type": "content", "title": "整体架构", "content": [
    "  前端技术栈（Vue 3）：",
    "  • Vite 5.0 - 构建工具",
    "  • Vue Router 4.2 - 路由管理",
    "  • Pinia 2.1 - 状态管理",
    "  • Element Plus 2.4 - UI组件库",
    "  • Axios 1.6 - HTTP请求",
    "",
    "  后端技术栈（Spring Boot 3.2）：",
    "  • Spring Boot - 核心框架",
    "  • MyBatis Plus 3.5.5 - ORM框架",
    "  • MySQL 8.0 - 数据库",
    "  • Knife4j - API文档",
    "  • Hutool - 工具类库",
    "  • Jsoup - 数据爬取",
    "  • Alipay SDK - 支付宝支付"
]})

slides_data.append({"type": "content", "title": "系统架构图", "content": [
    "",
    "                    ┌─────────────┐",
    "                    │   用户端     │",
    "                    │  Vue 3 +    │",
    "                    │ Element Plus │",
    "                    └──────┬──────┘",
    "                           │ HTTP/REST",
    "                    ┌──────▼──────┐",
    "                    │  Spring Boot │",
    "                    │   Backend    │",
    "                    │  Controller  │",
    "                    └──────┬──────┘",
    "                           │",
    "              ┌────────────┼────────────┐",
    "              │            │            │",
    "        ┌─────▼─────┐ ┌────▼────┐ ┌────▼────┐",
    "        │ MySQL DB  │ │  文件   │ │ 支付宝  │",
    "        │           │ │  存储   │ │  支付   │",
    "        └───────────┘ └─────────┘ └─────────┘",
    ""
]})

slides_data.append({"type": "content", "title": "核心技术亮点", "content": [
    "  亮点一：支付宝支付集成",
    "  • 集成支付宝沙箱环境，实现模拟支付功能",
    "  • 订单创建、支付回调、订单状态同步完整流程",
    "  • 安全可靠的支付接口封装",
    "",
    "  亮点二：爬虫自动获取资讯",
    "  • 使用Jsoup爬取考研资讯网站",
    "  • 自动解析HTML，提取标题、内容、发布时间",
    "  • 定时更新，保持资讯时效性",
    "",
    "  其他技术亮点：",
    "  • JWT Token身份认证",
    "  • AOP切面权限控制",
    "  • 统一异常处理"
]})

slides_data.append({"type": "section", "title": "四、数据库设计"})

slides_data.append({"type": "content", "title": "核心实体设计", "content": [
    "  用户相关：",
    "  • User - 用户信息表",
    "  • Address - 收货地址表",
    "",
    "  商城相关：",
    "  • Material - 学习资料表",
    "  • CartItem - 购物车表",
    "  • Order - 订单表",
    "  • OrderItem - 订单项表",
    "",
    "  内容相关：",
    "  • News - 资讯表",
    "  • Post - 帖子表",
    "  • Comment - 评论表",
    "  • Favorite - 收藏表",
    "  • Guide - 备考指南表",
    "",
    "  其他：",
    "  • Feedback - 反馈建议表"
]})

slides_data.append({"type": "content", "title": "用户与订单ER图", "content": [
    "  ER图关系说明：",
    "",
    "  User（用户） 1 ──n Address（地址）",
    "     │",
    "     └──n CartItem（购物车）",
    "     │",
    "     └──n Order（订单）──n OrderItem（订单项）",
    "                       │",
    "                       └─n Material（资料）",
    "",
    "  User（用户） 1 ──n Post（帖子）──n Comment（评论）",
    "     │",
    "     └──n Favorite（收藏）──n Material（资料）",
    "",
    "  关键设计：",
    "  • 采用自增主键标识每条记录",
    "  • 逻辑删除标记（deleted字段）",
    "  • 创建时间、更新时间自动管理",
    "  • 订单关联用户和收货地址"
]})

slides_data.append({"type": "section", "title": "五、功能模块展示"})

slides_data.append({"type": "content", "title": "用户模块", "content": [
    "  功能列表：",
    "  • 用户注册：输入用户名、密码、确认密码",
    "  • 用户登录：用户名密码验证，JWT Token返回",
    "  • 个人信息：查看和修改个人资料",
    "  • 收货地址：添加、编辑、删除收货地址",
    "  • 收藏管理：查看和管理已收藏的资料",
    "  • 反馈建议：提交问题反馈和建议",
    "",
    "  技术实现：",
    "  • 密码使用BCrypt加密存储",
    "  • JWT Token实现无状态认证",
    "  • 自定义注解实现权限控制"
]})

slides_data.append({"type": "content", "title": "资讯中心与爬虫模块", "content": [
    "  爬虫功能：",
    "  • 定时爬取考研资讯网站",
    "  • 自动解析HTML提取有效信息",
    "  • 去除广告和无效内容",
    "  • 数据清洗后入库",
    "",
    "  资讯展示：",
    "  • 资讯分类浏览",
    "  • 资讯详情展示",
    "  • 发布时间显示",
    "",
    "  核心代码：",
    "  Jsoup.connect(url).get() 获取页面",
    "  Element.select() 解析DOM",
    "  Document.text() 提取文本"
]})

slides_data.append({"type": "content", "title": "资料商城模块", "content": [
    "  商品浏览：",
    "  • 资料分类展示",
    "  • 搜索查找资料",
    "  • 资料详情页展示",
    "",
    "  购物车：",
    "  • 加入购物车",
    "  • 修改数量",
    "  • 删除商品",
    "",
    "  订单流程：",
    "  • 确认订单信息",
    "  • 选择收货地址",
    "  • 选择支付方式",
    "  • 支付宝支付",
    "  • 支付回调处理",
    "  • 订单完成"
]})

slides_data.append({"type": "content", "title": "交流论坛模块", "content": [
    "  帖子功能：",
    "  • 发布帖子（标题、内容）",
    "  • 编辑自己的帖子",
    "  • 删除自己的帖子",
    "",
    "  互动功能：",
    "  • 发表评论",
    "  • 点赞帖子",
    "  • 收藏帖子",
    "",
    "  浏览功能：",
    "  • 帖子列表浏览",
    "  • 帖子详情查看",
    "  • 查看评论列表",
    "",
    "  技术实现：",
    "  • 富文本编辑器支持",
    "  • 评论嵌套展示",
    "  • 点赞收藏状态管理"
]})

slides_data.append({"type": "content", "title": "管理后台模块", "content": [
    "  用户管理：",
    "  • 查看用户列表",
    "  • 禁用/启用用户",
    "",
    "  内容管理：",
    "  • 资讯管理（增删改）",
    "  • 帖子管理（审核、删除）",
    "  • 资料管理（上下架）",
    "  • 备考指南管理",
    "",
    "  订单管理：",
    "  • 查看所有订单",
    "  • 订单状态管理",
    "",
    "  反馈管理：",
    "  • 查看用户反馈",
    "  • 回复处理反馈"
]})

slides_data.append({"type": "section", "title": "六、系统测试与总结"})

slides_data.append({"type": "content", "title": "系统测试", "content": [
    "  功能测试：",
    "  • 用户注册登录功能测试 ✓",
    "  • 资料商城购物流程测试 ✓",
    "  • 支付宝支付功能测试 ✓",
    "  • 论坛发帖评论功能测试 ✓",
    "  • 管理后台功能测试 ✓",
    "",
    "  测试结果：",
    "  • 各功能模块运行正常",
    "  • 用户体验流畅",
    "  • 界面响应及时",
    "  • 安全机制有效"
]})

slides_data.append({"type": "content", "title": "项目总结", "content": [
    "  项目成果：",
    "  • 完成了一个功能完整的考研服务平台",
    "  • 实现了用户系统、商城、论坛、资讯等核心功能",
    "  • 集成了支付宝支付和爬虫资讯获取",
    "",
    "  创新点：",
    "  • 一站式服务平台设计",
    "  • 爬虫自动获取最新考研资讯",
    "  • 支付宝沙箱环境支付集成",
    "",
    "  不足与改进：",
    "  • 可增加实时聊天功能",
    "  • 可优化移动端适配",
    "  • 可扩展AI智能推荐功能"
]})

slides_data.append({"type": "title", "title": "感谢聆听", "subtitle": "\n\n欢迎各位老师批评指正！"})

for slide_data in slides_data:
    if slide_data["type"] == "title":
        add_title_slide(prs, slide_data["title"], slide_data.get("subtitle", ""))
    elif slide_data["type"] == "section":
        add_section_slide(prs, slide_data["title"])
    elif slide_data["type"] == "content":
        if slide_data.get("two_column"):
            add_content_slide(prs, slide_data["title"], slide_data["content"], two_column=True)
        else:
            add_content_slide(prs, slide_data["title"], slide_data["content"])

prs.save("研路无忧_毕业答辩.pptx")
print("PPT生成成功：研路无忧_毕业答辩.pptx")